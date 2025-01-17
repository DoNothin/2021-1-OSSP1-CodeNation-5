# -*- coding: utf-8 -*- 
from tika import parser
from pptx import Presentation
import olefile
import docx2txt
from krwordrank.hangle import normalize
from krwordrank.word import KRWordRank


# pdf to txt
def pdf_to_txt(path):
    pdf_path = path
    raw_pdf = parser.from_file(pdf_path) 
    result = raw_pdf['content'] 
    result = result.strip()
    #print(result)
    return result


# ppt to txt ---> 리스트에 txt 박스 단위로 반환 
def ppt_to_txt(path):
    ppt_path = Presentation(path)
    result = []
    for slide in ppt_path.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                result.append(paragraph.text)
    #print(result)
    return result


# hwp to txt ---> Prvtxt만 가능
# https://luji.tistory.com/18
def hwp_to_txt(path):
    hwp_file = olefile.OleFileIO(path)
    hwp_txt = hwp_file.openstream('Prvtext').read()
    result = hwp_txt.decode('UTF=16')
    #print (result)
    return result


# word to txt
def word_to_txt(path):
    result = docx2txt.process(path)
    #print(result)
    return result


# keyword extraction from txt
# https://lovit.github.io/nlp/2018/04/16/krwordrank/ 

def keyword_extraction(path):
    with open(path, 'r') as f:
        list_file = []
        for line in f:
            list_file.append(line)

    texts = list_file
    texts = [normalize(text, english=True, number=True) for text in texts]

    wordrank_extractor = KRWordRank(
        min_count = 5, # 단어의 최소 출현 빈도수 (그래프 생성 시)
        max_length = 10, # 단어의 최대 길이
        verbose = True
        )

    beta = 0.85    # PageRank의 decaying factor beta
    max_iter = 10

    keywords, rank, graph = wordrank_extractor.extract(texts, beta, max_iter)

    for word, r in sorted(keywords.items(), key=lambda x:x[1], reverse=True)[:5]: 
        print('%s' % (word))
        #print('%8s:\t%.4f' % (word, r))